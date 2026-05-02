"""
Assemble the final MATH GR5360 PowerPoint deck — fully Columbia-themed.

This builds the deck from scratch (uses the case-comp template only for
its slide-master dimensions) so every slide has consistent typography,
palette, and footer.

Output: report/presentation/MATH5360_Final_Group1.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Columbia palette
# ---------------------------------------------------------------------------
COL_NAVY = RGBColor(0x01, 0x21, 0x69)
COL_BLUE = RGBColor(0xB9, 0xD9, 0xEB)
COL_INK = RGBColor(0x1B, 0x36, 0x5D)
COL_GOLD = RGBColor(0xA2, 0x8D, 0x5B)
COL_RED = RGBColor(0xA0, 0x30, 0x33)
COL_GREEN = RGBColor(0x3F, 0x6F, 0x4A)
COL_GREY = RGBColor(0x9A, 0xA1, 0xA9)
COL_LIGHT = RGBColor(0xE6, 0xE9, 0xEE)
COL_CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)
COL_CREAM = RGBColor(0xF4, 0xF1, 0xEA)
COL_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = Path("/Users/nigelli/Downloads/5360-Presentation.pptx")
FIG_PRES = ROOT / "report" / "presentation" / "figures"
FIG_REPORT = ROOT / "report" / "figures"
OUT = ROOT / "report" / "presentation" / "MATH5360_Final_Group1.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def remove_all_slides(prs: Presentation) -> None:
    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def add_blank(prs: Presentation):
    layout = prs.slide_layouts[3]  # Title Only — minimal placeholders
    slide = prs.slides.add_slide(layout)
    # Hide any inherited title placeholder
    for shape in list(slide.shapes):
        if shape.has_text_frame and shape.is_placeholder:
            sp = shape._element  # noqa: SLF001
            sp.getparent().remove(sp)
    return slide


def paint_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text: str, *, size: float = 12,
             bold: bool = False, color: RGBColor = COL_INK,
             align=PP_ALIGN.LEFT, font: str = "Calibri", italic: bool = False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, *, size: float = 12,
                color: RGBColor = COL_INK, bullet: str = "•",
                bullet_color: RGBColor = COL_GOLD, line_spacing: float = 1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # Bullet
        rb = p.add_run()
        rb.text = f"{bullet}  "
        rb.font.size = Pt(size)
        rb.font.color.rgb = bullet_color
        rb.font.bold = True
        rb.font.name = "Calibri"
        # Body
        rt = p.add_run()
        rt.text = line
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
        rt.font.name = "Calibri"
    return box


def add_rect(slide, left, top, width, height, fill: RGBColor, *, line: RGBColor | None = None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_title_strip(slide, prs: Presentation, title: str, section: str = ""):
    sw = prs.slide_width
    add_text(slide, Inches(0.4), Inches(0.25), sw - Inches(0.8), Inches(0.55),
             title, size=24, bold=True, color=COL_NAVY)
    line = slide.shapes.add_connector(1, Inches(0.4), Inches(0.85), sw - Inches(0.4), Inches(0.85))
    line.line.color.rgb = COL_NAVY
    line.line.width = Pt(1.0)
    if section:
        add_text(slide, sw - Inches(2.6), Inches(0.25), Inches(2.2), Inches(0.4),
                 section.upper(), size=10, bold=True, color=COL_GOLD, align=PP_ALIGN.RIGHT)


def add_footer(slide, prs: Presentation, slide_no: int, total: int):
    sw, sh = prs.slide_width, prs.slide_height
    add_text(slide, Inches(0.4), sh - Inches(0.45), sw - Inches(0.8), Inches(0.3),
             "MATH GR5360 Final Project · Group 1 · Columbia MAFN — Spring 2026",
             size=9, color=COL_GREY)
    add_text(slide, sw - Inches(1.4), sh - Inches(0.45), Inches(1.0), Inches(0.3),
             f"{slide_no} / {total}", size=9, color=COL_GREY, align=PP_ALIGN.RIGHT)


def fit_image_in_box(slide, path: Path, box_left, box_top, box_w, box_h):
    if not path.exists():
        print(f"[warn] missing image: {path}")
        return None
    pic = slide.shapes.add_picture(str(path), box_left, box_top)
    iw, ih = pic.width, pic.height
    scale = min(box_w / iw, box_h / ih)
    nw = int(iw * scale)
    nh = int(ih * scale)
    pic.width = nw
    pic.height = nh
    pic.left = int(box_left + (box_w - nw) / 2)
    pic.top = int(box_top + (box_h - nh) / 2)
    return pic


# ---------------------------------------------------------------------------
# Slide builders — front matter (1-13)
# ---------------------------------------------------------------------------
def slide_title(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs)
    paint_bg(slide, COL_NAVY)
    sw, sh = prs.slide_width, prs.slide_height

    # Top thin gold line
    add_rect(slide, Inches(0.0), Inches(0.0), sw, Inches(0.05), COL_GOLD)

    # Pre-title small label
    add_text(slide, Inches(0.6), Inches(1.1), sw - Inches(1.2), Inches(0.4),
             "MATH GR5360  ·  MATHEMATICAL METHODS IN FINANCIAL PRICE ANALYSIS",
             size=11, bold=True, color=COL_GOLD)

    # Main title (two lines)
    add_text(slide, Inches(0.6), Inches(1.55), sw - Inches(1.2), Inches(1.4),
             "Trend-Following Across Markets",
             size=44, bold=True, color=COL_WHITE)
    add_text(slide, Inches(0.6), Inches(2.45), sw - Inches(1.2), Inches(1.0),
             "TY vs BTC via Channel WithDDControl",
             size=26, bold=False, color=COL_BLUE)

    # Subtitle
    add_text(slide, Inches(0.6), Inches(3.55), sw - Inches(1.2), Inches(0.5),
             "Walk-forward validation of horizon-specific trend inefficiencies",
             size=14, color=COL_LIGHT, italic=True)

    # Stat strip — 4 boxes
    strip_y = Inches(4.5)
    strip_h = Inches(1.4)
    strip_w = (sw - Inches(1.4)) / 4
    stats = [
        ("$68 336", "TY OOS net profit"),
        ("4.31×", "TY return on account"),
        ("$536 397", "BTC OOS net profit"),
        ("4.07×", "BTC return on account"),
    ]
    for i, (big, small) in enumerate(stats):
        x = Inches(0.7) + i * strip_w
        add_rect(slide, int(x), strip_y, int(strip_w - Inches(0.1)), strip_h,
                 COL_INK, line=COL_GOLD)
        add_text(slide, int(x), strip_y + Inches(0.22), int(strip_w - Inches(0.1)), Inches(0.6),
                 big, size=22, bold=True, color=COL_GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, int(x), strip_y + Inches(0.85), int(strip_w - Inches(0.1)), Inches(0.4),
                 small, size=10, color=COL_LIGHT, align=PP_ALIGN.CENTER)

    # Authors line
    add_text(slide, Inches(0.6), sh - Inches(1.0), sw - Inches(1.2), Inches(0.4),
             "Group 1   ·   Columbia MAFN   ·   Spring 2026",
             size=12, bold=True, color=COL_GOLD)

    # Footer
    add_text(slide, Inches(0.6), sh - Inches(0.55), sw - Inches(1.2), Inches(0.3),
             f"{slide_no} / {total}", size=9, color=COL_LIGHT, align=PP_ALIGN.RIGHT)


def slide_executive_summary(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Executive Summary", section="Overview")

    # Big quote
    add_text(slide, Inches(0.5), Inches(1.05), sw - Inches(1.0), Inches(0.7),
             '"Trend-following is real in both markets — but the usable horizon is market-specific."',
             size=18, bold=True, color=COL_NAVY, italic=True)

    # 3 cards: Thesis / Evidence / Results
    card_y = Inches(2.0)
    card_h = Inches(3.6)
    card_w = (sw - Inches(1.2)) / 3 - Inches(0.1)
    card_data = [
        ("THESIS",
         "TY exhibits weak / mixed behavior at intraday horizons but trends at multi-week scales. "
         "BTC mean-reverts at 1–4 day horizons but trends from ~12 days. Same Channel WithDDControl "
         "system, calibrated to each market's autocorrelation.",
         COL_NAVY),
        ("EVIDENCE",
         "Variance Ratio profile and Push-Response diagrams confirm trend at the multi-week horizon "
         "for TY (Spearman ρ = +0.59 at 18 sessions) and at the multi-day horizon for BTC "
         "(ρ = +0.67 at 12 days). Engines validated to float-64 parity.",
         COL_GOLD),
        ("RESULTS",
         "TY OOS: $68 336 / RoA 4.31× / Sharpe 0.31 across 1987–2026 (155 quarters). "
         "BTC OOS: $536 397 / RoA 4.07× / Sharpe 3.01 across 2023–2026 (7 quarters). "
         "Both ≈ 4× return-on-account, full slippage charged.",
         COL_RED),
    ]
    for i, (head, body, accent) in enumerate(card_data):
        x = Inches(0.5) + i * (card_w + Inches(0.15))
        add_rect(slide, int(x), card_y, int(card_w), card_h, COL_CREAM, line=accent)
        add_rect(slide, int(x), card_y, int(card_w), Inches(0.4), accent)
        add_text(slide, int(x) + Inches(0.2), card_y + Inches(0.05), int(card_w) - Inches(0.4),
                 Inches(0.35), head, size=12, bold=True, color=COL_WHITE)
        add_text(slide, int(x) + Inches(0.2), card_y + Inches(0.55),
                 int(card_w) - Inches(0.4), card_h - Inches(0.7),
                 body, size=11, color=COL_INK)

    # Bottom strip — key finding
    add_rect(slide, Inches(0.5), sh - Inches(1.4), sw - Inches(1.0), Inches(0.7), COL_NAVY)
    add_text(slide, Inches(0.7), sh - Inches(1.27), sw - Inches(1.4), Inches(0.5),
             "Key finding: One framework, calibrated to each market's autocorrelation structure, "
             "delivers ≈ 4× return-on-account out-of-sample on both TY and BTC.",
             size=12, bold=True, color=COL_WHITE)

    add_footer(slide, prs, slide_no, total)


def slide_roadmap(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Roadmap", section="Agenda")

    sections = [
        ("01", "Markets & Methodology",
         "Contract specifications, data audit, why TY × BTC pair, assignment scope."),
        ("02", "Diagnostics: Random Walk Tests",
         "Variance Ratio (Lo–MacKinlay) + Push-Response — locate the inefficiency in time."),
        ("03", "Strategy: Channel WithDDControl",
         "Breakout system + drawdown control. Walk-forward 4 yr IS / 1 quarter OOS, full grid."),
        ("04", "Out-of-Sample Performance",
         "Equity curves, drawdown family, trade ledgers, best/worst trade autopsies."),
        ("05", "Robustness",
         "IS-vs-OOS decay, parameter stability, T × τ sweep, Python ↔ C++ parity."),
        ("06", "Conclusion",
         "Same system, different inefficiency time-scales. Both markets earn ≈ 4× RoA OOS."),
    ]
    item_h = Inches(0.85)
    item_y0 = Inches(1.15)
    for i, (num, head, body) in enumerate(sections):
        y = item_y0 + i * (item_h + Inches(0.05))
        # Number column
        add_rect(slide, Inches(0.5), y, Inches(0.85), item_h, COL_NAVY)
        add_text(slide, Inches(0.5), y + Inches(0.18), Inches(0.85), Inches(0.5),
                 num, size=22, bold=True, color=COL_GOLD, align=PP_ALIGN.CENTER)
        # Body
        add_rect(slide, Inches(1.4), y, sw - Inches(1.9), item_h, COL_CREAM)
        add_text(slide, Inches(1.55), y + Inches(0.1), sw - Inches(2.2), Inches(0.4),
                 head, size=14, bold=True, color=COL_NAVY)
        add_text(slide, Inches(1.55), y + Inches(0.45), sw - Inches(2.2), Inches(0.4),
                 body, size=10.5, color=COL_INK)

    add_footer(slide, prs, slide_no, total)


def slide_assignment(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Assignment Brief", section="Methodology")

    items_left = [
        ("Variance Ratio test", "Confirm departure from random walk; identify horizon bands."),
        ("Push-Response test", "Measure how long a price impact persists after a large move."),
        ("Channel WithDDControl", "Port the main.m / ezread.m breakout system to a fast language."),
    ]
    items_right = [
        ("Walk-Forward Validation",
         "T = 4 yr IS, τ = 1 Q OOS; full grid over (L, S); objective Net Profit / Max Drawdown."),
        ("T × τ Sensitivity",
         "Robustness across IS = 1–10 yr and OOS = 1–4 Q (or 1–12 m)."),
        ("Two-market study",
         "Primary: TY (10-yr Treasury futures). Secondary: BTC (CME Bitcoin futures)."),
    ]

    col_w = (sw - Inches(1.3)) / 2
    for col_idx, items in enumerate([items_left, items_right]):
        x = Inches(0.5) + col_idx * (col_w + Inches(0.3))
        for i, (head, body) in enumerate(items):
            y = Inches(1.3) + i * Inches(1.5)
            # Number bullet
            add_rect(slide, x, y, Inches(0.35), Inches(0.35), COL_GOLD)
            add_text(slide, x, y + Inches(0.03), Inches(0.35), Inches(0.3),
                     "✓", size=14, bold=True, color=COL_WHITE, align=PP_ALIGN.CENTER)
            # Head
            add_text(slide, x + Inches(0.5), y, col_w - Inches(0.5), Inches(0.4),
                     head, size=14, bold=True, color=COL_NAVY)
            # Body
            add_text(slide, x + Inches(0.5), y + Inches(0.4), col_w - Inches(0.5), Inches(1.0),
                     body, size=11, color=COL_INK)

    # Bottom strip
    add_rect(slide, Inches(0.5), sh - Inches(1.3), sw - Inches(1.0), Inches(0.65), COL_CREAM)
    add_text(slide, Inches(0.7), sh - Inches(1.18), sw - Inches(1.4), Inches(0.4),
             "Grading rubric: \"judged on how close your results are to the expected ones\" — "
             "we report the 4× return-on-account both markets deliver under TF-Data slippage.",
             size=11, bold=False, color=COL_INK, italic=True)

    add_footer(slide, prs, slide_no, total)


def slide_market_specs(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Markets — contract specifications", section="Data & Markets")

    # Subtitle
    add_text(slide, Inches(0.5), Inches(1.0), sw - Inches(1.0), Inches(0.4),
             "Same system, different inefficiency time-scales — calibration follows the contract.",
             size=12, italic=True, color=COL_GREY)

    # Build the table manually as rectangles for full styling control
    headers = ["Field", "TY — 10-yr Treasury futures", "BTC — CME Bitcoin futures"]
    rows = [
        ("Exchange", "CBOT / CME", "CME"),
        ("Currency", "USD", "USD"),
        ("Contract size", "$100 000 face value", "5 BTC"),
        ("Point value", "$1 000", "$5"),
        ("Tick size", "1/64 of a point (0.015625)", "5.00 points"),
        ("Tick value", "$15.625", "$25.00"),
        ("Slippage (TF Data)", "$18.625 / round-turn", "$25.00 / round-turn"),
        ("Settlement", "Physical delivery", "Cash-settled (BRR index)"),
        ("Trading hours (local)", "07:20 – 14:00 (Chicago)", "17:00 – 16:00 (Globex, ~23 h)"),
        ("Bars per session", "80", "276"),
        ("Sample window", "03 Jan 1983 – 10 Apr 2026", "18 Dec 2017 – 10 Apr 2026"),
        ("5-minute bars", "863 887", "590 436"),
    ]

    table_x = Inches(0.5)
    table_y = Inches(1.5)
    table_w = sw - Inches(1.0)
    col_w = [Inches(2.6), (table_w - Inches(2.6)) / 2, (table_w - Inches(2.6)) / 2]
    row_h = Inches(0.32)

    # Header
    cx = table_x
    for j, h in enumerate(headers):
        add_rect(slide, cx, table_y, col_w[j], row_h + Inches(0.05), COL_NAVY)
        add_text(slide, cx + Inches(0.1), table_y + Inches(0.04), col_w[j] - Inches(0.2),
                 row_h, h, size=11, bold=True, color=COL_WHITE)
        cx += col_w[j]

    # Body rows
    for i, row in enumerate(rows):
        y = table_y + row_h + Inches(0.05) + i * row_h
        # Zebra
        if i % 2 == 0:
            add_rect(slide, table_x, y, table_w, row_h, COL_CREAM)
        cx = table_x
        for j, val in enumerate(row):
            color = COL_NAVY if j == 0 else COL_INK
            bold = j == 0
            add_text(slide, cx + Inches(0.1), y + Inches(0.04), col_w[j] - Inches(0.2),
                     row_h, val, size=10.5, bold=bold, color=color)
            cx += col_w[j]

    # Bottom note
    add_text(slide, Inches(0.5), sh - Inches(0.95), sw - Inches(1.0), Inches(0.4),
             "Slippage from TF Data column V; point value from column H. Sessions cross-checked "
             "against Bloomberg DES screens. BTC's 23-hour Globex session means 276 bars/day vs TY's 80.",
             size=9.5, color=COL_GREY, italic=True)

    add_footer(slide, prs, slide_no, total)


def slide_market_overview(prs: Presentation, slide_no: int, total: int, market: str):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    if market == "TY":
        title = "Primary market — TY (10-yr Treasury futures)"
        section = "Markets"
        img = FIG_PRES / "front_ty_price_history.png"
        bullets = [
            "43-year continuous backadjusted history — the longest dataset in the project.",
            "Macro-driven contract: Fed cycles, inflation prints, fiscal news.",
            "Deep, liquid CBOT order book; tight spreads; small intraday slippage.",
            "Annualised vol ≈ 4–6%; intraday signal-to-noise low.",
            "Trend horizon emerges at multi-week scales (push-response peaks at 18 sessions).",
        ]
    else:
        title = "Secondary market — BTC (CME Bitcoin futures)"
        section = "Markets"
        img = FIG_PRES / "front_btc_price_history.png"
        bullets = [
            "8.4-year history from contract inception (Dec 2017).",
            "Sentiment-driven flows: retail, algorithms, regulatory news, exchange events.",
            "Thinner order book; large flows leave persistent imprints.",
            "Annualised vol ≈ 30–60%; intraday signal-to-noise high.",
            "Trend horizon emerges at ≈ 12 days; mean-reverts on intraday timescales.",
        ]
    add_title_strip(slide, prs, title, section=section)

    # Image takes ~70% width on left, bullets right
    img_w = int(sw * 0.62)
    fit_image_in_box(slide, img, Inches(0.4), Inches(1.05), img_w, Inches(4.7))
    # Bullets
    bx = Inches(0.4) + img_w + Inches(0.2)
    bw = sw - bx - Inches(0.4)
    add_text(slide, bx, Inches(1.05), bw, Inches(0.4),
             "Market profile", size=13, bold=True, color=COL_NAVY)
    add_bullets(slide, bx, Inches(1.5), bw, Inches(4.5), bullets, size=11)

    # Caption
    if market == "TY":
        cap = ("TY's price action over 1983–2026 is the canonical macro-bond curve: secular highs "
               "into 2020–21, then the 2022–24 hiking cycle. Channel breakouts at 18-session scale capture each leg.")
    else:
        cap = ("BTC's 2017–2026 price chart is dominated by the 2021 cycle, the 2022 winter, and "
               "the 2024–25 cycle that powered 7/7 winning OOS quarters in our walk-forward.")
    add_text(slide, Inches(0.5), sh - Inches(0.95), sw - Inches(1.0), Inches(0.45),
             cap, size=10.5, color=COL_INK)

    add_footer(slide, prs, slide_no, total)


def slide_horizons_context(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Why the same trend principle uses different horizons",
                    section="Markets")
    fit_image_in_box(slide, FIG_PRES / "front_horizons_context.png",
                     Inches(0.4), Inches(1.0), sw - Inches(0.8), sh - Inches(2.0))
    add_text(slide, Inches(0.5), sh - Inches(0.95), sw - Inches(1.0), Inches(0.45),
             "Same channel-breakout principle on both markets — the optimiser self-selects "
             "L ≈ 1 920 bars (~24 sessions) for TY and L ≈ 276 bars (~1 day) for BTC, exactly the "
             "horizons where push-response is positive on each market.",
             size=10.5, color=COL_INK, italic=True)
    add_footer(slide, prs, slide_no, total)


def slide_vr_definition(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Random Walk Test 1 — Variance Ratio", section="Diagnostics")
    fit_image_in_box(slide, FIG_PRES / "front_vr_test_card.png",
                     Inches(0.4), Inches(1.0), sw - Inches(0.8), sh - Inches(2.0))
    add_text(slide, Inches(0.5), sh - Inches(0.95), sw - Inches(1.0), Inches(0.45),
             "We compute VR(q) on price differences over a logarithmic grid of horizons up to 40 "
             "sessions for TY and 20 days for BTC, with the heteroskedasticity-robust Z₂* statistic.",
             size=10.5, color=COL_INK)
    add_footer(slide, prs, slide_no, total)


def slide_pr_definition(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Random Walk Test 2 — Push-Response", section="Diagnostics")
    fit_image_in_box(slide, FIG_PRES / "front_pr_test_card.png",
                     Inches(0.4), Inches(1.0), sw - Inches(0.8), sh - Inches(2.0))
    add_text(slide, Inches(0.5), sh - Inches(0.95), sw - Inches(1.0), Inches(0.45),
             "We bin pushes into deciles, plot the conditional mean of the response per decile, "
             "and report the Spearman ρ over the binned (push, response) pairs.",
             size=10.5, color=COL_INK)
    add_footer(slide, prs, slide_no, total)


def slide_central_thesis(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Central Thesis", section="Diagnostics")

    # Big quote
    add_text(slide, Inches(0.5), Inches(1.05), sw - Inches(1.0), Inches(1.2),
             "TY shows weak / mixed short-horizon behaviour but trends at multi-week scales.\n"
             "BTC mean-reverts at 1–4 day horizons but trends from ~12 days.\n\n"
             "→ the same Channel WithDDControl framework calibrates to each market.",
             size=15, bold=True, color=COL_NAVY)

    # 3 supporting cards
    card_y = Inches(3.6)
    card_h = Inches(2.4)
    card_w = (sw - Inches(1.4)) / 3 - Inches(0.05)
    cards = [
        ("01", "Different autocorrelation",
         "BTC's variance-ratio profile dips deeper (0.82 at ~9 days) than TY's (0.89 at ~10 sessions). "
         "Push-response is positive at ~12 days for BTC and ~18 sessions for TY. The two markets "
         "literally live in different time-scales of the trend / mean-revert spectrum."),
        ("02", "Same framework",
         "Channel WithDDControl picks (L, S) by walk-forward optimisation on Net Profit / Max Drawdown. "
         "Because the trend signal is at a shorter horizon for BTC, the optimiser self-selects a "
         "much smaller L there (276) than on TY (1 920) — no heuristic, just the data."),
        ("03", "OOS validation",
         "TY OOS: $68 336 / RoA 4.31× / Sharpe 0.31. BTC OOS: $536 397 / RoA 4.07× / Sharpe 3.01. "
         "Both markets earn ≈ 4× return-on-account out-of-sample, with full Bloomberg/TF-Data "
         "slippage charged on every round-turn."),
    ]
    for i, (num, head, body) in enumerate(cards):
        x = Inches(0.5) + i * (card_w + Inches(0.1))
        # Number stripe
        add_rect(slide, int(x), card_y, int(card_w), Inches(0.5), COL_NAVY)
        add_text(slide, int(x) + Inches(0.2), card_y + Inches(0.07), Inches(1.0), Inches(0.4),
                 num, size=18, bold=True, color=COL_GOLD)
        add_text(slide, int(x) + Inches(0.9), card_y + Inches(0.13), int(card_w) - Inches(1.1),
                 Inches(0.35), head, size=12, bold=True, color=COL_WHITE)
        # Body
        add_rect(slide, int(x), card_y + Inches(0.5), int(card_w), card_h - Inches(0.5),
                 COL_CREAM)
        add_text(slide, int(x) + Inches(0.2), card_y + Inches(0.62),
                 int(card_w) - Inches(0.4), card_h - Inches(0.65),
                 body, size=10.5, color=COL_INK)

    add_footer(slide, prs, slide_no, total)


def slide_horizon_spectrum(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Horizon spectrum — where each market lives",
                    section="Diagnostics")
    fit_image_in_box(slide, FIG_PRES / "front_horizon_spectrum.png",
                     Inches(0.4), Inches(1.0), sw - Inches(0.8), Inches(4.2))

    # Side-by-side annotations below the figure
    box_y = Inches(5.4)
    box_h = Inches(1.3)
    box_w = (sw - Inches(1.0)) / 2 - Inches(0.1)

    # TY box
    add_rect(slide, Inches(0.5), box_y, box_w, box_h, COL_CREAM, line=COL_NAVY)
    add_text(slide, Inches(0.7), box_y + Inches(0.1), box_w - Inches(0.4), Inches(0.4),
             "TY  ·  18-session trend regime", size=13, bold=True, color=COL_NAVY)
    add_text(slide, Inches(0.7), box_y + Inches(0.5), box_w - Inches(0.4), Inches(0.7),
             "Push-response Spearman ρ = +0.59 at q = 1 440 bars (≈ 18 sessions). "
             "Optimiser picks modal L* ≈ 1 920 bars (≈ 24 sessions), S* ≈ 1%.",
             size=10.5, color=COL_INK)

    # BTC box
    add_rect(slide, Inches(0.7) + box_w, box_y, box_w, box_h, COL_CREAM, line=COL_GOLD)
    add_text(slide, Inches(0.9) + box_w, box_y + Inches(0.1), box_w - Inches(0.4), Inches(0.4),
             "BTC  ·  12-day trend regime", size=13, bold=True, color=COL_GOLD)
    add_text(slide, Inches(0.9) + box_w, box_y + Inches(0.5), box_w - Inches(0.4), Inches(0.7),
             "Push-response Spearman ρ = +0.67 at q = 3 456 bars (≈ 12 days). "
             "Optimiser picks modal L* ≈ 276 bars (≈ 1 day) or 1 104 bars (≈ 4 days), S* ≈ 1%.",
             size=10.5, color=COL_INK)

    add_footer(slide, prs, slide_no, total)


# ---------------------------------------------------------------------------
# Slide builders — strategy & walk-forward (16, 17)
# ---------------------------------------------------------------------------
def slide_strategy_mechanics(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Strategy — Channel WithDDControl", section="Strategy")

    # Two columns: parameter table left, logic flow right
    col_w = (sw - Inches(1.3)) / 2

    # Left column — parameters
    px = Inches(0.5)
    add_text(slide, px, Inches(1.05), col_w, Inches(0.4),
             "Core parameters", size=14, bold=True, color=COL_NAVY)

    params = [
        ("L", "Channel length",
         "Number of bars in the rolling High/Low channel. Entry triggers when Close breaks "
         "above High(L−1) (long) or below Low(L−1) (short). Captures the trending horizon."),
        ("S", "Drawdown stop",
         "Fraction of running peak equity. Position closes when the running drawdown exceeds "
         "S × peak. The risk control that makes the breakout system survive losing streaks."),
        ("Slippage", "Round-turn cost",
         "Debited on every position change. $18.625 for TY, $25.00 for BTC, per TF Data column V. "
         "Charged in full on every entry and exit."),
    ]
    py = Inches(1.5)
    for tag, head, body in params:
        # Tag chip
        add_rect(slide, px, py, Inches(0.65), Inches(0.65), COL_NAVY)
        add_text(slide, px, py + Inches(0.13), Inches(0.65), Inches(0.45),
                 tag, size=18, bold=True, color=COL_GOLD, align=PP_ALIGN.CENTER)
        # Head + body
        add_text(slide, px + Inches(0.85), py, col_w - Inches(0.85), Inches(0.3),
                 head, size=12, bold=True, color=COL_NAVY)
        add_text(slide, px + Inches(0.85), py + Inches(0.32), col_w - Inches(0.85), Inches(1.4),
                 body, size=10.5, color=COL_INK)
        py += Inches(1.55)

    # Right column — logic flow
    qx = Inches(0.7) + col_w
    add_text(slide, qx, Inches(1.05), col_w, Inches(0.4),
             "Logic flow per bar", size=14, bold=True, color=COL_NAVY)
    flow = [
        ("1", "Compute rolling channel",
         "high_band = max(High[t−L .. t−1])\nlow_band  = min(Low[t−L .. t−1])"),
        ("2", "Check breakout",
         "If Close[t] > high_band → enter LONG\nIf Close[t] < low_band → enter SHORT"),
        ("3", "Track equity & drawdown",
         "Mark-to-market each bar. Track running peak and drawdown."),
        ("4", "Apply drawdown stop",
         "If (peak − equity) ≥ S × peak → close position; charge slippage."),
        ("5", "Emit closed-trade ledger",
         "Record entry/exit time, direction, prices, duration, PnL, slippage."),
    ]
    qy = Inches(1.5)
    for num, head, body in flow:
        add_rect(slide, qx, qy, Inches(0.4), Inches(0.4), COL_GOLD)
        add_text(slide, qx, qy + Inches(0.05), Inches(0.4), Inches(0.32),
                 num, size=12, bold=True, color=COL_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, qx + Inches(0.55), qy, col_w - Inches(0.55), Inches(0.3),
                 head, size=11.5, bold=True, color=COL_NAVY)
        add_text(slide, qx + Inches(0.55), qy + Inches(0.3),
                 col_w - Inches(0.55), Inches(0.7),
                 body, size=10, color=COL_INK,
                 font="DejaVu Sans Mono")
        qy += Inches(0.95)

    add_footer(slide, prs, slide_no, total)


def slide_walkforward(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Walk-forward design", section="Strategy")
    fit_image_in_box(slide, FIG_PRES / "front_walkforward_schematic.png",
                     Inches(0.4), Inches(1.0), sw - Inches(0.8), Inches(3.6))
    # Three boxes with knobs
    boxes = [
        ("T = 4 years", "In-sample window",
         "Chosen so each IS spans multiple regimes for TY (rate cycles) and the full BTC inception."),
        ("τ = 1 quarter", "Out-of-sample window",
         "Short enough to test near-term generalisation; long enough to gather statistically meaningful trades."),
        ("Net P / Max DD", "Optimisation objective",
         "Calmar-style: rewards capital efficiency, penalises drawdowns. Tie-break on smaller L, then smaller S."),
    ]
    bx = Inches(0.5)
    by = Inches(4.8)
    bw = (sw - Inches(1.2)) / 3 - Inches(0.05)
    bh = Inches(1.5)
    for i, (k, title, body) in enumerate(boxes):
        x = bx + i * (bw + Inches(0.1))
        add_rect(slide, x, by, bw, Inches(0.45), COL_NAVY)
        add_text(slide, x + Inches(0.2), by + Inches(0.08), bw - Inches(0.4),
                 Inches(0.35), k, size=12, bold=True, color=COL_GOLD)
        add_rect(slide, x, by + Inches(0.45), bw, bh - Inches(0.45), COL_CREAM)
        add_text(slide, x + Inches(0.2), by + Inches(0.55), bw - Inches(0.4), Inches(0.3),
                 title, size=10.5, bold=True, color=COL_NAVY)
        add_text(slide, x + Inches(0.2), by + Inches(0.85), bw - Inches(0.4), Inches(1.0),
                 body, size=10, color=COL_INK)

    add_footer(slide, prs, slide_no, total)


# ---------------------------------------------------------------------------
# Existing-style slide builders (results 18-35) — kept from prior version
# ---------------------------------------------------------------------------
def slide_image_full(prs: Presentation, *, title: str, section: str,
                     image: Path, caption: str, slide_no: int, total: int,
                     side_bullets: list[str] | None = None):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, title, section)
    if side_bullets:
        img_w = int(sw * 0.62)
        fit_image_in_box(slide, image, Inches(0.4), Inches(1.0), img_w, sh - Inches(2.0))
        rx = Inches(0.4) + img_w + Inches(0.2)
        rw = sw - rx - Inches(0.4)
        add_bullets(slide, rx, Inches(1.1), rw, sh - Inches(2.2), side_bullets, size=11)
    else:
        fit_image_in_box(slide, image, Inches(0.4), Inches(1.0), sw - Inches(0.8), sh - Inches(2.0))
    add_text(slide, Inches(0.4), sh - Inches(0.95), sw - Inches(0.8), Inches(0.45),
             caption, size=10.5, color=COL_INK)
    add_footer(slide, prs, slide_no, total)


def slide_section_divider(prs: Presentation, *, label: str, big: str,
                          slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_NAVY)
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, Inches(0.0), Inches(0.0), sw, Inches(0.05), COL_GOLD)
    add_text(slide, Inches(0.6), Inches(2.6), sw - Inches(1.2), Inches(0.5),
             label.upper(), size=14, bold=True, color=COL_GOLD)
    add_text(slide, Inches(0.6), Inches(3.1), sw - Inches(1.2), Inches(2.0),
             big, size=36, bold=True, color=COL_WHITE)
    add_text(slide, Inches(0.6), sh - Inches(0.45), sw - Inches(1.2), Inches(0.3),
             "MATH GR5360 Final Project · Group 1 · Columbia MAFN — Spring 2026",
             size=9, color=COL_LIGHT)
    add_text(slide, sw - Inches(1.4), sh - Inches(0.45), Inches(1.0), Inches(0.3),
             f"{slide_no} / {total}", size=9, color=COL_LIGHT, align=PP_ALIGN.RIGHT)


def slide_metrics_table(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Performance metrics", section="Results")
    fit_image_in_box(slide, FIG_PRES / "slide_01_performance_metrics.png",
                     Inches(0.4), Inches(1.0), sw - Inches(0.8), sh - Inches(2.0))
    add_text(slide, Inches(0.4), sh - Inches(0.95), sw - Inches(0.8), Inches(0.45),
             "Both markets earn ≈ 4× Return on Account out-of-sample. TY's signal is slow and "
             "low-vol; BTC's is fast and high-vol — same Channel WithDDControl framework, "
             "calibrated to each market's autocorrelation structure.",
             size=10.5, color=COL_INK)
    add_footer(slide, prs, slide_no, total)


def slide_conclusion(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Conclusion", section="Summary")
    add_text(slide, Inches(0.6), Inches(1.1), sw - Inches(1.2), Inches(0.7),
             '"Trend-following is real in both markets — but the usable horizon is market-specific."',
             size=18, bold=True, color=COL_NAVY, italic=True)
    cards = [
        ("01", "Diagnostics",
         "Variance-ratio test and push-response confirm trend-following at the multi-week horizon "
         "for TY (Spearman ρ ≈ 0.59 at 18 sessions) and at the multi-day horizon for BTC "
         "(ρ ≈ 0.67 at 12 days). Same statistical framework — different frequency bands."),
        ("02", "Strategy",
         "Channel WithDDControl ported from main.m / ezread.m to Python (Numba JIT) and C++17. "
         "Walk-forward 4-yr IS / 1-Q OOS over the full ChnLen × StpPct grid, "
         "objective = Net Profit / Max Drawdown."),
        ("03", "OOS Results",
         "TY: $68 336 net / 4.31× RoA / Sharpe 0.31 over 1987–2026 (155 quarters). "
         "BTC: $536 397 net / 4.07× RoA / Sharpe 3.01 over 2023–2026 (7 quarters). "
         "Both deliver ≈ 4× return-on-account with full slippage charged."),
        ("04", "Robustness",
         "Python ↔ C++ parity to float-64 on every metric and every closed trade. "
         "T × τ sweep confirms (4 yr, 1 Q) is the well-behaved choice. 1-minute resolution "
         "earns marginally more (RoA 4.61×) but does not change the picture."),
    ]
    card_y = Inches(2.0)
    card_h = Inches(1.0)
    for i, (num, head, body) in enumerate(cards):
        y = card_y + i * (card_h + Inches(0.05))
        add_text(slide, Inches(0.6), y, Inches(0.7), card_h, num,
                 size=28, bold=True, color=COL_GOLD)
        add_text(slide, Inches(1.4), y, Inches(2.1), Inches(0.4), head,
                 size=13, bold=True, color=COL_NAVY)
        add_text(slide, Inches(3.6), y, sw - Inches(4.0), card_h, body,
                 size=10.5, color=COL_INK)
    add_footer(slide, prs, slide_no, total)


def slide_appendix(prs: Presentation, slide_no: int, total: int):
    slide = add_blank(prs); paint_bg(slide, COL_WHITE)
    sw, sh = prs.slide_width, prs.slide_height
    add_title_strip(slide, prs, "Appendix — repository, parity, key figures",
                    section="Appendix")
    items = [
        "Repository: github.com/nl2992/MATH5360_Final_Project",
        "Walk-forward Python artifacts: results/walkforward/{TY_5m, BTC_5m, TY_1m}/",
        "C++ parity artifacts: results/cpp_parity/{TY, BTC}/",
        "Diagnostics (VR + Push-Response cached tables): results/diagnostics/",
        "Comprehensive write-up: report/FINAL_REPORT.md",
        "Slide-companion narrative: report/presentation/demo.md",
        "All figures regenerated by: scripts/build_final_report_figures.py, "
        "build_presentation_figures.py and build_front_matter_figures.py",
        "Python ↔ C++ parity: every metric and every closed trade match to float-64 — "
        "see results/walkforward/python_cpp_fidelity_comparison.csv",
    ]
    add_bullets(slide, Inches(0.6), Inches(1.2), sw - Inches(1.2), sh - Inches(2.2),
                items, size=12)
    add_text(slide, Inches(0.6), sh - Inches(1.0), sw - Inches(1.2), Inches(0.5),
             "Channel WithDDControl ported from main.m / ezread.m. Python core JIT-compiled with "
             "Numba; C++17 reference engine confirms parity to float-64.",
             size=10, color=COL_GREY, italic=True)
    add_footer(slide, prs, slide_no, total)


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------
def main() -> None:
    prs = Presentation(str(TEMPLATE))
    print(f"start: {len(prs.slides)} slides")
    remove_all_slides(prs)
    print(f"after wipe: {len(prs.slides)} slides")

    TOTAL = 41

    # ----- Front matter (1-13) -----
    slide_title(prs, 1, TOTAL)
    slide_executive_summary(prs, 2, TOTAL)
    slide_roadmap(prs, 3, TOTAL)
    slide_assignment(prs, 4, TOTAL)
    slide_market_specs(prs, 5, TOTAL)
    slide_market_overview(prs, 6, TOTAL, market="TY")
    slide_market_overview(prs, 7, TOTAL, market="BTC")
    slide_horizons_context(prs, 8, TOTAL)
    slide_vr_definition(prs, 9, TOTAL)
    slide_pr_definition(prs, 10, TOTAL)
    slide_central_thesis(prs, 11, TOTAL)
    slide_horizon_spectrum(prs, 12, TOTAL)
    slide_section_divider(prs, label="Diagnostics",
                          big="Random walk tests on TY and BTC — VR + Push-Response across horizons and decades",
                          slide_no=13, total=TOTAL)

    # ----- Diagnostic figures (14-19) -----
    # 14 — TY price + implied yield (replicates the reference)
    slide_image_full(
        prs,
        title="TY price and implied 10-yr yield, 1983–2026",
        section="Diagnostics",
        image=FIG_PRES / "repl_ty_price_yield.png",
        caption=(
            "Top: TY 10-yr Treasury futures daily close. Bottom: implied 10-yr yield, recovered "
            "from the futures price by inverting a 6%-coupon 20-period bond pricing equation. "
            "Reproduces the structural decline in yields from ~14% (1984) to ~0.5% (2020) before "
            "the 2022–24 hiking cycle returned yields to ~4–5%."
        ),
        slide_no=14, total=TOTAL,
    )
    # 15 — VR vs q (single-curve, q to 5000 — replicates the reference)
    slide_image_full(
        prs,
        title="Variance Ratio vs q — TY 5-min price differences",
        section="Diagnostics",
        image=FIG_PRES / "repl_ty_vr_curve.png",
        caption=(
            "Single-window VR(q) on TY 5-min price differences out to q = 5 000 bars (~62.5 trading days). "
            "VR dips to ≈ 0.89 around q = 800 (~10 sessions) and recovers toward 0.96 around q = 3 500 (~44 sessions). "
            "Across the entire profile VR < 1, but the gap closes at multi-week horizons — exactly where the "
            "push-response test will show positive ρ."
        ),
        slide_no=15, total=TOTAL,
    )
    # 16 — VR by past 10/20/30/40 year lookbacks (replicates the reference)
    slide_image_full(
        prs,
        title="Variance Ratio Test for TY: past 10/20/30/40 years",
        section="Diagnostics",
        image=FIG_PRES / "repl_ty_vr_lookback.png",
        caption=(
            "How the VR profile depends on the lookback window. The Past-10-year curve crosses above 1 "
            "around q ≈ 4 000 — Treasuries have been *trending* in the 2016–2026 sample. Longer lookbacks "
            "smooth this into the long-run profile, but the trend signal at multi-week horizons is preserved."
        ),
        slide_no=16, total=TOTAL,
    )
    # 17 — VR by backward 10-year windows (replicates the reference)
    slide_image_full(
        prs,
        title="Variance Ratio Test by backward 10-year windows — TY",
        section="Diagnostics",
        image=FIG_PRES / "repl_ty_vr_decade_windows.png",
        caption=(
            "Decade-by-decade decomposition. The 2016–2026 decade clearly trends at long q (VR > 1 above "
            "q ≈ 3 000). The 2006–2016 decade is the most mean-reverting (the QE/ZIRP era). The earlier "
            "windows sit closer to the random-walk null — regime change is real and visible in the data."
        ),
        slide_no=17, total=TOTAL,
    )
    # 18 — Push-response grid TY (replicates the reference)
    slide_image_full(
        prs,
        title="Push-Response Test for TY: short → medium horizons",
        section="Diagnostics",
        image=FIG_PRES / "repl_ty_pr_grid.png",
        caption=(
            "Conditional mean of the forward response, per push-bin, at twelve τ values from 1 to 350 bars. "
            "At τ = 1 the response is mean-reverting (negative slope). Around τ = 32–96 the slope flattens, "
            "and from τ = 144 onwards the response acquires a clear negative slope on extreme pushes — "
            "the multi-week trend regime."
        ),
        slide_no=18, total=TOTAL,
    )
    # 19 — Push-response combined diagram (the report figure)
    slide_image_full(
        prs,
        title="Evidence — Push-Response at the optimal horizon (TY vs BTC)",
        section="Diagnostics",
        image=FIG_REPORT / "fig_push_response.png",
        caption=(
            "Conditional mean response per push-decile with bin-level standard errors at the *optimal* horizon "
            "for each market. TY at q = 1 440 bars (~18 sessions): Spearman ρ ≈ +0.59 (p ≈ 0.06). "
            "BTC at q = 3 456 bars (~12 days): ρ ≈ +0.67 (p ≈ 0.02). These horizons drive the walk-forward picks of L*."
        ),
        slide_no=19, total=TOTAL,
    )

    # ----- Strategy & walk-forward (20-21) -----
    slide_strategy_mechanics(prs, 20, TOTAL)
    slide_walkforward(prs, 21, TOTAL)

    # ----- Performance section (22-37) -----
    slide_section_divider(prs, label="Out-of-Sample Performance",
                          big="Walk-forward results for both markets",
                          slide_no=22, total=TOTAL)
    slide_metrics_table(prs, 23, TOTAL)
    slide_image_full(
        prs,
        title="TY — out-of-sample equity & portfolio position",
        section="Primary Market — TY",
        image=FIG_PRES / "slide_02_ty_equity_position.png",
        caption=(
            "Channel WithDDControl on TY (10-yr Treasury futures) over 1987-06 → 2026-03 — "
            "155 quarterly walk-forward windows stitched, $100k initial equity. "
            "Net OOS profit $68 336 / Max DD $15 865 / Return on Account 4.31× / Sharpe 0.31 / 395 trades. "
            "Position panel shows long/short/flat state at bar resolution."
        ),
        slide_no=24, total=TOTAL,
    )
    slide_image_full(
        prs,
        title="TY — Chekhlov drawdown family",
        section="Primary Market — TY",
        image=FIG_PRES / "slide_03_ty_drawdown_family.png",
        caption=(
            "Top: % off running peak. Bottom: $ off running peak. Max DD ≈ 11% / $15.9k. "
            "CDD(α=0.05) ≈ $13.3k indicates the worst 5% of drawdown bars are concentrated near the max. "
            "Long underwater stretches are structural for a 1.5%/yr trend-following bond strategy."
        ),
        slide_no=25, total=TOTAL,
    )
    # NEW — % return distribution for TY (replaces the dollar-PnL hist on the deck)
    slide_image_full(
        prs,
        title="TY — out-of-sample trade % return distribution",
        section="Primary Market — TY",
        image=FIG_PRES / "repl_ty_pct_returns.png",
        caption=(
            "Per-trade % returns — left panel = price-move % captured per trade; right panel = "
            "equity-return % per trade (PnL / $100k). Mean equity-return per trade is −0.18% (small "
            "losing tail of −2.95%); winners reach +8.17% on the COVID flight-to-safety long."
        ),
        slide_no=26, total=TOTAL,
    )
    slide_image_full(
        prs,
        title="TY — out-of-sample dollar-PnL distribution",
        section="Primary Market — TY",
        image=FIG_PRES / "slide_04_ty_trade_distribution.png",
        caption=(
            "Same trades shown in dollar terms. Win rate 33%, average winner $1 265 vs average loser −$897. "
            "The dollar view exaggerates outliers because the position size is fixed at one contract — "
            "the % view above is the cleaner cross-market comparator."
        ),
        slide_no=27, total=TOTAL,
    )
    slide_image_full(
        prs,
        title="TY — most profitable OOS trade (autopsy)",
        section="Primary Market — TY",
        image=FIG_PRES / "slide_05_ty_best_trade.png",
        caption=(
            "21 Jan 2020 LONG 129.56 → 137.75 over 50 days (+$8 170). "
            "Channel breakout above the 1 920-bar high right as COVID drove a flight-to-safety bond rally. "
            "Trailing-equity stop fired after the initial impulse decayed — textbook trend payoff."
        ),
        slide_no=28, total=TOTAL,
    )
    slide_image_full(
        prs,
        title="TY — worst OOS trade · why it got cooked",
        section="Primary Market — TY",
        image=FIG_PRES / "slide_06_ty_worst_trade.png",
        side_bullets=[
            "22 Feb 2002 LONG @ 107.38 — broke above the 3 200-bar (~40-day) high.",
            "Treasuries reversed almost immediately on hawkish-Fed signals at the late-Feb FOMC.",
            "Price slid 2.93 pts in 12 days; trailing stop fired at $-2 952.",
            "Classic 'breakout caught at the local top' before mean-reversion fired.",
            "Wide L = 3 200 channel made the per-bar stop physically distant — loss compounded before exit.",
            "The push-response diagnostic for short horizons is near zero; some breakouts get faded.",
        ],
        caption="Source: Group 1 walk-forward (TF Data 5-min OHLC, $100k initial equity)",
        slide_no=29, total=TOTAL,
    )
    slide_image_full(
        prs,
        title="TY — walk-forward parameter stability",
        section="Primary Market — TY",
        image=FIG_PRES / "slide_07_ty_param_stability.png",
        caption=(
            "Channel length L converges to a tight cluster around 1 920 bars (≈ 24 trading days). "
            "Drawdown stop S almost always at 1% of running equity. Optimiser does not flip wildly — "
            "the chosen objective (Net Profit / Max Drawdown) is well-behaved on TY."
        ),
        slide_no=30, total=TOTAL,
    )

    slide_section_divider(prs, label="Secondary Market — BTC",
                          big="Same framework, faster horizon",
                          slide_no=31, total=TOTAL)

    slide_image_full(
        prs,
        title="BTC — out-of-sample equity & portfolio position",
        section="Secondary Market — BTC",
        image=FIG_PRES / "slide_02_btc_equity_position.png",
        caption=(
            "Channel WithDDControl on BTC (CME Bitcoin futures) over 2023-08 → 2026-02 — 7 quarterly "
            "OOS windows. Net OOS profit $536 397 / Max DD $131 729 / Return on Account 4.07× / "
            "Sharpe 3.01 / 1 094 trades. Equity grows from $100k to $636k driven by the 2024-25 cycle."
        ),
        slide_no=32, total=TOTAL,
    )
    slide_image_full(
        prs,
        title="BTC — Chekhlov drawdown family",
        section="Secondary Market — BTC",
        image=FIG_PRES / "slide_03_btc_drawdown_family.png",
        caption=(
            "Max DD ≈ 22% / $131.7k. Recovery is fast — most underwater stretches are weeks, not "
            "months. CDD(α=0.05) ≈ $111.8k. Larger absolute dollar swings than TY, but realised "
            "vol is ~8× higher so the % drawdown footprint is comparable."
        ),
        slide_no=33, total=TOTAL,
    )
    # NEW — % return distribution for BTC
    slide_image_full(
        prs,
        title="BTC — out-of-sample trade % return distribution",
        section="Secondary Market — BTC",
        image=FIG_PRES / "repl_btc_pct_returns.png",
        caption=(
            "Per-trade % returns — left panel = price-move % captured; right panel = equity-return %. "
            "BTC's right tail is dramatic: best equity-return trade is +45% (the 25-minute weekend-gap "
            "long); worst is −10.6% (the channel-break short that BTC pumped through). Mean is +0.49%."
        ),
        slide_no=34, total=TOTAL,
    )
    slide_image_full(
        prs,
        title="BTC — out-of-sample dollar-PnL distribution",
        section="Secondary Market — BTC",
        image=FIG_PRES / "slide_04_btc_trade_distribution.png",
        caption=(
            "Win rate 42% with average winner $4 327 vs average loser −$2 285 — same right-skew as TY "
            "but with a higher hit rate. Profit factor 1.37. The intraday push-response is positive "
            "at the right horizons, supporting the higher hit rate."
        ),
        slide_no=35, total=TOTAL,
    )
    slide_image_full(
        prs,
        title="BTC — most profitable OOS trade (autopsy)",
        section="Secondary Market — BTC",
        image=FIG_PRES / "slide_05_btc_best_trade.png",
        caption=(
            "02 Mar 2025 LONG 85 720 → 94 748 in 25 minutes (+$45 115). "
            "Channel break above the 276-bar (1-day) high; weekend gap aligned with the breakout direction. "
            "$9 028 price move × $5 BTC point value − $25 round-turn slippage = +$45 115 in 5 bars."
        ),
        slide_no=36, total=TOTAL,
    )
    slide_image_full(
        prs,
        title="BTC — worst OOS trade · why it got cooked",
        section="Secondary Market — BTC",
        image=FIG_PRES / "slide_06_btc_worst_trade.png",
        side_bullets=[
            "22 Aug 2025 SHORT @ 112 075 — broke below the 276-bar (1-day) low.",
            "Within 90 minutes BTC pumped $2 115 (+1.9%) against the position.",
            "Tight 1% drawdown stop fired at the second bar after the move ($114 190).",
            "Channel breakout in BTC's mean-reverting 1-day regime — recurring failure mode.",
            "The push-response diagram already flagged BTC as mean-reverting at the 1-day horizon "
            "(Spearman ρ = −0.38).",
            "The −$10 600 loss is ≈ 8% of the OOS Max DD — paid once to stay in the longer-horizon trend regime.",
        ],
        caption="Source: Group 1 walk-forward (TF Data 5-min OHLC, $100k initial equity)",
        slide_no=37, total=TOTAL,
    )
    slide_image_full(
        prs,
        title="BTC — walk-forward parameter stability",
        section="Secondary Market — BTC",
        image=FIG_PRES / "slide_07_btc_param_stability.png",
        caption=(
            "BTC has only 7 quarterly OOS windows (the sample's IS warmup eats ~4 years from inception). "
            "Optimiser cycles between L = 276 (1 day) in choppy regimes and L = 1 104 (4 days) in the "
            "late-2025 trend phase. S* fixed at 1%. Instability is structural — BTC's regime is shifting."
        ),
        slide_no=38, total=TOTAL,
    )

    # ----- Robustness + close (39-41) -----
    slide_image_full(
        prs,
        title="In-sample vs out-of-sample decay",
        section="Robustness",
        image=FIG_REPORT / "fig_is_oos_metrics.png",
        caption=(
            "Sharpe decay: TY 0.41 → 0.31 (-24%); BTC 4.52 → 3.01 (-33%). "
            "Profit-per-quarter decay: TY 26%, BTC 130% (BTC's IS objective favoured smaller L which "
            "truncated the most explosive late-2024 trends — sample-specific, not a forward claim). "
            "Quarterly OOS hit rate: TY 54.8% (85/155), BTC 100% (7/7)."
        ),
        slide_no=39, total=TOTAL,
    )
    slide_conclusion(prs, 40, TOTAL)
    slide_appendix(prs, 41, TOTAL)

    print(f"final: {len(prs.slides)} slides")
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"[ok] saved -> {OUT}")


if __name__ == "__main__":
    main()
